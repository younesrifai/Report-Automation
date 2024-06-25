import yfinance as yf
import pandas as pd
import urllib.request
from datetime import datetime
from pptx import Presentation

tickers = pd.read_csv("tickers.csv", header = None)[0]

stock_variations = []

for ticker in tickers:
    x = yf.Ticker(ticker)
    hist = x.history(period='1d')
    hist.reset_index(drop=True)
    var = (hist['Close'].values - hist['Open'].values)/hist['Open'].values
    try:
        stock_variations.append([ticker, var[0]])
    except: 
        stock_variations.append([ticker, var])

df = pd.DataFrame(stock_variations, columns=['Ticker', 'Variation'])
df = df.sort_values(by='Variation', ascending=False)

best_stock = df.iloc[0]['Ticker']
best_var = str(round(df.iloc[0]['Variation']*100,2)) + '%'
worst_stock = df.iloc[-1]['Ticker']
worst_var = str(round(df.iloc[-1]['Variation']*100,2)) + '%'

best_news = yf.Ticker(best_stock).news[0]
best_title = best_news['title']
if len(best_title) > 95:
    best_title = best_title[0:90] + '...'
best_link = best_news['link']
best_publisher = best_news['publisher']

worst_news = yf.Ticker(worst_stock).news[0]
worst_title = worst_news['title']
if len(worst_title) > 95:
    worst_title = worst_title[0:90] + '...'
worst_link = worst_news['link']
worst_publisher = worst_news['publisher']

today = datetime.today().strftime('%Y-%m-%d')
try:
    best_image_url = yf.Ticker(best_stock).news[0]['thumbnail']['resolutions'][0]['url']
    urllib.request.urlretrieve(best_image_url, 'images/' + str(today) + " best.png") 
    best_img_path = 'images/' + str(today) +  " best.png"
except:
    try:
        best_image_url = yf.Ticker(best_stock).news[1]['thumbnail']['resolutions'][0]['url']
        urllib.request.urlretrieve(best_image_url, 'images/' + str(today) + " best.png") 
        best_img_path = 'images/' + str(today) +  " best.png"
    except:
        try:
            best_image_url = yf.Ticker(best_stock).news[2]['thumbnail']['resolutions'][0]['url']
            urllib.request.urlretrieve(best_image_url, 'images/' + str(today) + " best.png") 
            best_img_path = 'images/' + str(today) +  " best.png"
        except:
            best_img_path = 'images/best generic.jpg'

try:
    worst_image_url = yf.Ticker(worst_stock).news[0]['thumbnail']['resolutions'][0]['url']
    urllib.request.urlretrieve(worst_image_url, 'images/' + str(today) + " worst.png")
    worst_img_path = 'images/' + str(today) + " worst.png"
except:
    worst_image_path = 'images/best generic.jpg'

template_path = 'templates/Movers_Template.pptx'
presentation = Presentation(template_path)

data = [
    {
        'Date': today,
        'Best Stock' :  yf.Ticker(best_stock).info['longName'],
        'Best Ticker' : best_stock,
        'Best Variation' : best_var,
        'Best Headline' : best_title,
        'Best Image' : best_img_path,
        'Best Source' : best_publisher,
        'Worst Stock' : yf.Ticker(worst_stock).info['longName'],
        'Worst Ticker' : worst_stock,
        'Worst Variation' : worst_var,
        'Worst Headline' : worst_title,
        'Worst Image' : worst_img_path,
        'Worst Source' : worst_publisher,
    }
]
print(data)

for slide, x in zip(presentation.slides, data):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if '{Date}' in run.text:
                    run.text = run.text.replace('{Date}', x['Date'])
                if '{Best Stock}' in run.text:
                    run.text = run.text.replace('{Best Stock}', x['Best Stock'])
                if '{Best Ticker}' in run.text:
                    run.text = run.text.replace('{Best Ticker}', x['Best Ticker'])
                if '{Best Variation}' in run.text:
                    run.text = run.text.replace('{Best Variation}', x['Best Variation'])
                if '{Best Headline}' in run.text:
                    run.text = run.text.replace('{Best Headline}', x['Best Headline'])
                if '{Best Source}' in run.text:
                    run.text = run.text.replace('{Best Source}', x['Best Source'])
                if '{Worst Stock}' in run.text:
                    run.text = run.text.replace('{Worst Stock}', x['Worst Stock'])
                if '{Worst Ticker}' in run.text:
                    run.text = run.text.replace('{Worst Ticker}', x['Worst Ticker'])
                if '{Worst Variation}' in run.text:
                    run.text = run.text.replace('{Worst Variation}', x['Worst Variation'])
                if '{Worst Headline}' in run.text:
                    run.text = run.text.replace('{Worst Headline}', x['Worst Headline'])
                if '{Worst Source}' in run.text:
                    run.text = run.text.replace('{Worst Source}', x['Worst Source'])

    for shape in slide.shapes:
        if shape.name in x:
            if shape.name == 'Best Image':
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                best_image_path = x['Best Image']
                slide.shapes.add_picture(best_image_path, left, top, width, height)
            
                # Remove the original placeholder shape
                spTree = slide.shapes._spTree
                spTree.remove(shape._element)

            if shape.name == 'Worst Image':
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                worst_image_path = x['Worst Image']
                slide.shapes.add_picture(worst_image_path, left, top, width, height)
            
                # Remove the original placeholder shape
                spTree = slide.shapes._spTree
                spTree.remove(shape._element)

presentation.save(f'Stock Movement - {today}.pptx')